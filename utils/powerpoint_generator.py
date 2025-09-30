try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. PowerPoint export will not be available.")

from datetime import datetime

class PowerpointGenerator:
    def __init__(self, topic):
        self.topic = topic

    def create_title_slide(self, prs):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = self.topic
        subtitle.text = f"Research Presentation\nGenerated: {datetime.now().strftime('%Y-%m-%d')}"

    def create_content_slide(self, prs, title_text, bullets):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content_box = slide.placeholders[1]

        title.text = title_text

        text_frame = content_box.text_frame
        text_frame.clear()

        for bullet in bullets[:5]:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

    def generate_powerpoint(self, slides_data, output_path="research_presentation.pptx"):
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            self.create_title_slide(prs)

            for slide_data in slides_data:
                self.create_content_slide(prs, slide_data["title"], slide_data["bullets"])

            # Save directly to file
            prs.save(output_path)
            return output_path
            
        except Exception as e:
            print(f"Error generating PowerPoint: {str(e)}")
            raise
