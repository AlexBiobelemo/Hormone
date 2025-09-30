import streamlit as st
import google.generativeai as genai
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
from datetime import datetime

# Page configuration
st.set_page_config(
    page_title="Hormonal Profile Test Research Assistant",
    page_icon="🧬",
    layout="wide"
)

# Initialize session state
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'api_key' not in st.session_state:
    st.session_state.api_key = None

# Get API key from secrets or user input
def get_api_key():
    # Try to get from Streamlit secrets first (for deployment)
    try:
        return st.secrets["GEMINI_API_KEY"]
    except (KeyError, FileNotFoundError):
        # Fall back to session state or user input
        return st.session_state.api_key

# Sidebar for API key
with st.sidebar:
    st.title("⚙️ Configuration")
    
    # Check if API key exists in secrets
    secret_key = None
    try:
        secret_key = st.secrets["GEMINI_API_KEY"]
        st.success("✅ API Key loaded from secrets!")
        st.session_state.api_key = secret_key
        genai.configure(api_key=secret_key)
    except (KeyError, FileNotFoundError):
        st.info("💡 No API key found in secrets. Please enter manually.")
        api_key = st.text_input("Enter Google Gemini API Key", type="password", value=st.session_state.api_key or "")
        if api_key:
            st.session_state.api_key = api_key
            genai.configure(api_key=api_key)
            st.success("API Key configured!")
    
    st.markdown("---")
    st.markdown("### About")
    st.info("This application provides expert assistance on Hormonal Profile Tests, helping with research, generating detailed reports in PDF format, and creating professional PowerPoint presentations.")
    
    st.markdown("---")
    st.markdown("### 🔒 Security Note")
    st.markdown("""
    **For Streamlit Cloud Deployment:**
    1. Go to your app settings
    2. Navigate to Secrets
    3. Add your API key as:
    ```
    GEMINI_API_KEY = "your-api-key-here"
    ```
    """)

# System prompt for the AI
SYSTEM_PROMPT = """You are an expert medical research assistant specializing in Hormonal Profile Tests and endocrinology. Your knowledge includes:

**Core Expertise:**
- Hormonal analysis and interpretation of test results
- Reproductive health hormones (Estrogen, Progesterone, FSH, LH, Testosterone, Prolactin)
- Thyroid hormones (TSH, T3, T4) and their metabolic impact
- Sex-specific hormonal profiles for men and women
- Conditions like PCOS, infertility, menopause, and thyroid disorders
- Clinical testing procedures and interpretation

**Key Hormones You Specialize In:**

For Women:
- Estrogen (Estradiol E2): Female sex hormone, menstrual cycle regulation
- Progesterone: Pregnancy preparation, ovulation confirmation
- FSH: Ovarian follicle stimulation, menopause indicators
- LH: Ovulation trigger, reproductive function
- Testosterone: Sex drive, bone health, PCOS association
- Prolactin: Milk production, testosterone synthesis influence
- Thyroid hormones: Metabolic and reproductive impact

For Men:
- Testosterone: Male sexual characteristics, sperm production, libido
- DHT: Potent testosterone form, male development
- FSH: Sperm production regulation
- LH: Testosterone secretion stimulation
- Prolactin: Testosterone synthesis influence

**Testing Knowledge:**
- Blood sample collection procedures
- Timing considerations (especially menstrual cycle phases)
- Normal ranges by age and sex
- Interpretation of results and abnormalities
- Clinical indications: infertility, menstrual irregularities, PCOS, thyroid disorders, hormone therapy monitoring

Provide detailed, accurate, evidence-based information. Use clear medical terminology while remaining accessible. Cite clinical relevance and practical applications. Always emphasize consulting healthcare professionals for diagnosis and treatment."""

# Main title
st.title("🧬 Hormonal Profile Test Research Assistant")
st.markdown("### Expert guidance on hormonal analysis, testing, and interpretation")

# Create tabs
tab1, tab2, tab3 = st.tabs(["💬 Research Chat", "📄 PDF Generation", "📊 PowerPoint Generation"])

# Tab 1: Research Chat
with tab1:
    st.header("Research Assistant")
    st.markdown("Ask questions about hormonal profile tests, specific hormones, testing procedures, or interpretation of results.")
    
    if not st.session_state.api_key:
        st.warning("⚠️ Please enter your Google Gemini API key in the sidebar to start chatting.")
    else:
        # Display chat history
        for message in st.session_state.chat_history:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
        
        # Chat input
        if prompt := st.chat_input("Ask about hormonal profile tests..."):
            # Add user message to chat history
            st.session_state.chat_history.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)
            
            # Generate response
            with st.chat_message("assistant"):
                with st.spinner("Researching..."):
                    try:
                        model = genai.GenerativeModel('gemini-2.5-flash')
                        full_prompt = f"{SYSTEM_PROMPT}\n\nUser Question: {prompt}"
                        response = model.generate_content(full_prompt)
                        st.markdown(response.text)
                        st.session_state.chat_history.append({"role": "assistant", "content": response.text})
                    except Exception as e:
                        st.error(f"Error: {str(e)}")
        
        # Clear chat button
        if st.button("🗑️ Clear Chat History"):
            st.session_state.chat_history = []
            st.rerun()

# Tab 2: PDF Generation
with tab2:
    st.header("Generate Research PDF")
    st.markdown("Create a comprehensive PDF report on hormonal profile testing topics.")
    
    if not st.session_state.api_key:
        st.warning("⚠️ Please enter your Google Gemini API key in the sidebar to generate PDFs.")
    else:
        col1, col2 = st.columns([2, 1])
        
        with col1:
            pdf_topic = st.text_input("Enter research topic:", placeholder="e.g., FSH levels in menopause diagnosis")
            
        with col2:
            detail_level = st.selectbox("Detail Level:", ["Comprehensive", "Standard", "Summary"])
        
        if st.button("📄 Generate PDF Report", type="primary"):
            if pdf_topic:
                with st.spinner("Generating comprehensive research report..."):
                    try:
                        model = genai.GenerativeModel('gemini-2.5-flash')
                        
                        # Generate content based on detail level
                        detail_instructions = {
                            "Comprehensive": "Provide an extensive, detailed analysis with clinical data, research findings, and practical applications. Include sections on: Introduction, Clinical Significance, Normal Ranges, Testing Procedures, Interpretation Guidelines, Associated Conditions, and Clinical Recommendations.",
                            "Standard": "Provide a thorough overview with key clinical information, testing details, and interpretation guidelines. Include sections on: Overview, Key Hormones, Testing Process, Normal Values, and Clinical Applications.",
                            "Summary": "Provide a concise summary with essential information about the topic. Include: Key Points, Testing Basics, and Clinical Relevance."
                        }
                        
                        prompt = f"{SYSTEM_PROMPT}\n\nCreate a detailed research report on: {pdf_topic}\n\n{detail_instructions[detail_level]}\n\nFormat the response with clear sections and subsections."
                        
                        response = model.generate_content(prompt)
                        content = response.text
                        
                        # Create PDF
                        pdf = FPDF()
                        pdf.add_page()
                        pdf.set_auto_page_break(auto=True, margin=15)
                        
                        # Title
                        pdf.set_font("Arial", "B", 16)
                        pdf.cell(0, 10, "Hormonal Profile Test Research Report", ln=True, align="C")
                        pdf.ln(5)
                        
                        # Topic
                        pdf.set_font("Arial", "B", 14)
                        pdf.multi_cell(0, 10, f"Topic: {pdf_topic}")
                        pdf.ln(2)
                        
                        # Metadata
                        pdf.set_font("Arial", "I", 10)
                        pdf.cell(0, 8, f"Generated: {datetime.now().strftime('%B %d, %Y at %H:%M')}", ln=True)
                        pdf.cell(0, 8, f"Detail Level: {detail_level}", ln=True)
                        pdf.ln(5)
                        
                        # Content
                        pdf.set_font("Arial", "", 11)
                        
                        # Process content line by line
                        for line in content.split('\n'):
                            if line.strip():
                                if line.startswith('#'):
                                    # Headers
                                    pdf.ln(3)
                                    pdf.set_font("Arial", "B", 12)
                                    pdf.multi_cell(0, 8, line.replace('#', '').strip())
                                    pdf.set_font("Arial", "", 11)
                                    pdf.ln(2)
                                elif line.startswith('*') or line.startswith('-'):
                                    # Bullet points
                                    pdf.multi_cell(0, 6, f"  • {line[1:].strip()}")
                                else:
                                    # Regular text
                                    pdf.multi_cell(0, 6, line)
                            else:
                                pdf.ln(3)
                        
                        # Footer
                        pdf.ln(10)
                        pdf.set_font("Arial", "I", 9)
                        pdf.multi_cell(0, 5, "Disclaimer: This report is for educational and research purposes only. Always consult qualified healthcare professionals for medical advice, diagnosis, or treatment.")
                        
                        # Save PDF to bytes
                        pdf_output = pdf.output(dest='S').encode('latin-1')
                        
                        st.success("✅ PDF generated successfully!")
                        st.download_button(
                            label="⬇️ Download PDF Report",
                            data=pdf_output,
                            file_name=f"hormone_research_{pdf_topic.replace(' ', '_')[:30]}_{datetime.now().strftime('%Y%m%d')}.pdf",
                            mime="application/pdf"
                        )
                        
                        # Preview
                        with st.expander("📖 Preview Content"):
                            st.markdown(content)
                            
                    except Exception as e:
                        st.error(f"Error generating PDF: {str(e)}")
            else:
                st.warning("Please enter a research topic.")

# Tab 3: PowerPoint Generation
with tab3:
    st.header("Generate Research PowerPoint")
    st.markdown("Create a professional PowerPoint presentation on hormonal profile testing topics.")
    
    if not st.session_state.api_key:
        st.warning("⚠️ Please enter your Google Gemini API key in the sidebar to generate presentations.")
    else:
        col1, col2 = st.columns([2, 1])
        
        with col1:
            ppt_topic = st.text_input("Enter presentation topic:", placeholder="e.g., Understanding PCOS hormonal markers")
            
        with col2:
            num_slides = st.slider("Number of slides:", min_value=5, max_value=15, value=8)
        
        if st.button("📊 Generate PowerPoint Presentation", type="primary"):
            if ppt_topic:
                with st.spinner("Creating professional presentation..."):
                    try:
                        model = genai.GenerativeModel('gemini-2.5-flash')
                        
                        # Generate structured content for slides
                        prompt = f"""{SYSTEM_PROMPT}

Create a professional PowerPoint presentation outline on: {ppt_topic}

Generate exactly {num_slides} slides with the following structure:
- Slide 1: Title slide
- Slides 2-{num_slides-1}: Content slides
- Slide {num_slides}: Summary/Conclusion slide

For each slide, provide:
1. Slide Title (clear and concise)
2. 3-5 bullet points with detailed content
3. Key takeaway or clinical note

Format your response as:
SLIDE [number]: [Title]
- Bullet point 1
- Bullet point 2
- Bullet point 3
- Bullet point 4 (if applicable)
- Bullet point 5 (if applicable)

Ensure content is medically accurate, clinically relevant, and professionally presented."""
                        
                        response = model.generate_content(prompt)
                        content = response.text
                        
                        # Create PowerPoint
                        prs = Presentation()
                        prs.slide_width = Inches(10)
                        prs.slide_height = Inches(7.5)
                        
                        # Parse content into slides
                        slides_data = []
                        current_slide = None
                        
                        for line in content.split('\n'):
                            line = line.strip()
                            if line.startswith('SLIDE'):
                                if current_slide:
                                    slides_data.append(current_slide)
                                title = line.split(':', 1)[1].strip() if ':' in line else "Slide"
                                current_slide = {"title": title, "bullets": []}
                            elif line.startswith('-') or line.startswith('•'):
                                if current_slide:
                                    current_slide["bullets"].append(line[1:].strip())
                        
                        if current_slide:
                            slides_data.append(current_slide)
                        
                        # Create slides
                        for i, slide_data in enumerate(slides_data):
                            if i == 0:
                                # Title slide
                                slide_layout = prs.slide_layouts[0]
                                slide = prs.slides.add_slide(slide_layout)
                                title = slide.shapes.title
                                subtitle = slide.placeholders[1]
                                
                                title.text = ppt_topic
                                subtitle.text = f"Hormonal Profile Test Research\nGenerated: {datetime.now().strftime('%B %d, %Y')}"
                            else:
                                # Content slide
                                slide_layout = prs.slide_layouts[1]
                                slide = prs.slides.add_slide(slide_layout)
                                title = slide.shapes.title
                                content_box = slide.placeholders[1]
                                
                                title.text = slide_data["title"]
                                
                                text_frame = content_box.text_frame
                                text_frame.clear()
                                
                                for bullet in slide_data["bullets"][:5]:
                                    p = text_frame.add_paragraph()
                                    p.text = bullet
                                    p.level = 0
                                    p.font.size = Pt(18)
                        
                        # Save PowerPoint to bytes
                        ppt_bytes = io.BytesIO()
                        prs.save(ppt_bytes)
                        ppt_bytes.seek(0)
                        
                        st.success("✅ PowerPoint generated successfully!")
                        st.download_button(
                            label="⬇️ Download PowerPoint Presentation",
                            data=ppt_bytes,
                            file_name=f"hormone_presentation_{ppt_topic.replace(' ', '_')[:30]}_{datetime.now().strftime('%Y%m%d')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        
                        # Preview
                        with st.expander("📖 Preview Slide Outline"):
                            st.markdown(content)
                            
                    except Exception as e:
                        st.error(f"Error generating PowerPoint: {str(e)}")
            else:
                st.warning("Please enter a presentation topic.")

# Footer
st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #666; padding: 20px;'>
    <p><strong>Hormonal Profile Test Research Assistant</strong></p>
    <p>This application is designed for educational and research purposes only.</p>
    <p>Always consult qualified healthcare professionals for medical advice, diagnosis, or treatment.</p>
</div>
""", unsafe_allow_html=True)
